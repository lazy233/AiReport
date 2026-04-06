"""python-pptx：形状解析、幻灯片遍历、生成结果写回。"""
from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.text import MSO_UNDERLINE

from ppt_report.constants import PAGE_TYPE_LABELS
from ppt_report.services.chapter_ref_images import image_file_path
from ppt_report.services.chapter_reference_resolve import ppt_reference_slot_rows

CONNECTOR_SHAPE_TYPE = getattr(MSO_SHAPE_TYPE, "CONNECTOR", None)


def emu_to_cm(value: int) -> float:
    return round(value / 360000, 2)


def infer_role(component_type: str, text: str) -> str:
    text_l = (text or "").strip().lower()

    if component_type == "title":
        return "用于概括当前页的核心主题"
    if component_type == "subtitle":
        return "用于补充标题背景信息"
    if component_type == "text":
        if any(token in text_l for token in ["总结", "结论", "thanks", "thank", "结束"]):
            return "用于给出结论或收尾信息"
        if len(text_l) > 120:
            return "用于承载详细说明内容"
        return "用于补充说明和要点描述"
    if component_type == "table":
        return "用于结构化展示行列数据并便于对比"
    if component_type == "chart":
        return "用于可视化展示数据趋势、占比或对比关系"
    if component_type == "image":
        return "用于视觉说明、示例展示或背景强化"
    if component_type == "line_or_arrow":
        return "用于表达流程、连接关系或方向指引"
    if component_type == "group":
        return "用于将多个元素组合成一个语义单元"
    if component_type == "placeholder":
        return "用于预留内容区域，便于快速套用模板"
    if component_type == "body":
        return "用于正文段落或列表（占位符正文区）"
    if component_type == "table_cell":
        return "用于在表格单元格内展示数据，可与整表其它单元格对照阅读"
    return "用于页面排版或视觉组织"


def map_shape_type(shape) -> str:
    if shape.is_placeholder:
        ph_type = shape.placeholder_format.type
        if ph_type in (
            PP_PLACEHOLDER.TITLE,
            PP_PLACEHOLDER.CENTER_TITLE,
            PP_PLACEHOLDER.VERTICAL_TITLE,
        ):
            return "title"
        if ph_type == PP_PLACEHOLDER.SUBTITLE:
            return "subtitle"
        if ph_type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.VERTICAL_BODY):
            return "body"
        return "placeholder"

    if getattr(shape, "has_table", False):
        return "table"
    if getattr(shape, "has_chart", False):
        return "chart"

    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return "image"
    if shape.shape_type == MSO_SHAPE_TYPE.LINE or (
        CONNECTOR_SHAPE_TYPE is not None and shape.shape_type == CONNECTOR_SHAPE_TYPE
    ):
        return "line_or_arrow"
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        return "group"
    if getattr(shape, "has_text_frame", False):
        return "text"

    return str(shape.shape_type).lower()


def _max_run_font_pt(shape) -> float | None:
    if not getattr(shape, "has_text_frame", False):
        return None
    best: float | None = None
    try:
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                sz = r.font.size
                if sz is None:
                    continue
                pt = float(getattr(sz, "pt", 0) or 0)
                if pt <= 0:
                    continue
                if best is None or pt > best:
                    best = pt
    except (AttributeError, TypeError, ValueError):
        return best
    return best


def refine_heading_cap_category(shape, comp_type: str) -> str | None:
    if str(comp_type).lower() != "text":
        return None
    name = shape.name or ""
    name_l = name.lower()
    if any(x in name for x in ("副标题", "副题")) or "subtitle" in name_l:
        return "subtitle"
    if "副" in name and "标题" in name:
        return "subtitle"
    if any(x in name for x in ("主标题", "大标题")):
        return "title"
    if any(k in name_l for k in ("title", "heading")) and "sub" not in name_l:
        return "title"
    if "标题" in name:
        return "title"
    if any(x in name for x in ("小标题", "节标题", "章节标题")) or "section" in name_l:
        return "placeholder"

    h_cm = emu_to_cm(shape.height)
    w_cm = emu_to_cm(shape.width)
    pt = _max_run_font_pt(shape)
    if pt is not None and w_cm >= 3.5 and h_cm <= 3.8:
        if pt >= 26:
            return "title"
        if pt >= 16:
            return "subtitle"
    if 0 < h_cm <= 2.2 and w_cm >= 8:
        return "placeholder"
    return None


def _font_color_rgb_hex(font) -> str | None:
    """字体前景色为显式 RGB 时返回 #RRGGBB；主题色/未定义返回 None。"""
    try:
        c = font.color
        if c is None:
            return None
        if c.type != MSO_COLOR_TYPE.RGB:
            return None
        rgb = c.rgb
        if rgb is None:
            return None
        r, g, b = int(rgb[0]), int(rgb[1]), int(rgb[2])
        return f"#{r:02x}{g:02x}{b:02x}"
    except (AttributeError, IndexError, TypeError, ValueError):
        return None


def _serialize_run_style(run) -> dict[str, Any]:
    """将单个 text run 的字体样式转为可 JSON 序列化的字典（仅包含有明确值的字段）。"""
    font = run.font
    d: dict[str, Any] = {}
    if font.bold is not None:
        d["bold"] = bool(font.bold)
    if font.italic is not None:
        d["italic"] = bool(font.italic)
    u = font.underline
    if u is not None:
        d["underline"] = u != MSO_UNDERLINE.NONE
    sz = font.size
    if sz is not None:
        try:
            d["font_size_pt"] = round(float(sz.pt), 2)
        except (AttributeError, TypeError, ValueError):
            pass
    name = font.name
    if name and str(name).strip():
        d["font_name"] = str(name).strip()
    hx = _font_color_rgb_hex(font)
    if hx:
        d["color_rgb"] = hx
    return d


def _normalize_run_text_for_export(raw: str) -> str:
    """python-pptx 用 \\v 表示段内软换行（Shift+Enter / <a:br>），导出为 \\n 便于阅读与前端换行。"""
    return (raw or "").replace("\v", "\n")


def extract_text_frame_style_runs(text_frame) -> list[dict[str, Any]]:
    """
    按文档顺序输出文本 run 列表；段与段之间插入仅含换行符的片段（无样式字段）。
    段内软换行在 run 的 text 中为 \\n（由 \\v 归一化而来）。
    与 extract_text 的空白折叠策略无关，便于还原行与样式。
    """
    if text_frame is None:
        return []
    out: list[dict[str, Any]] = []
    paras = list(text_frame.paragraphs)
    for pi, para in enumerate(paras):
        for r in para.runs:
            row: dict[str, Any] = {"text": _normalize_run_text_for_export(r.text or "")}
            row.update(_serialize_run_style(r))
            out.append(row)
        if pi < len(paras) - 1:
            out.append({"text": "\n"})
    return out


def extract_text(shape) -> str:
    if getattr(shape, "has_text_frame", False):
        text = shape.text_frame.text or ""
        return " ".join(text.split())
    return ""


def extract_text_from_cell(cell) -> str:
    try:
        tf = getattr(cell, "text_frame", None)
        if tf is None:
            return ""
        text = tf.text or ""
        return " ".join(text.split())
    except (AttributeError, TypeError, ValueError):
        return ""


def flatten_table_text(shape) -> tuple[str, int, int]:
    """返回 (制表符分隔的整表文本, 行数, 列数)；非表格则为 ('', 0, 0)。"""
    if not getattr(shape, "has_table", False):
        return "", 0, 0
    try:
        tbl = shape.table
        nrows = len(tbl.rows)
        ncols = len(tbl.columns)
        lines: list[str] = []
        for r in range(nrows):
            cells: list[str] = []
            for c in range(ncols):
                cells.append(extract_text_from_cell(tbl.cell(r, c)))
            lines.append("\t".join(cells))
        return "\n".join(lines), nrows, ncols
    except (AttributeError, TypeError, ValueError):
        return "", 0, 0


def estimate_max_chars(width_cm: float, height_cm: float, component_type: str) -> int:
    if component_type == "title":
        return max(8, min(30, int(width_cm * 2.2)))
    if component_type == "subtitle":
        return max(12, min(50, int(width_cm * 3)))
    area = max(1.0, width_cm * height_cm)
    return max(15, min(1200, int(area * 8)))


def parse_shape_component(shape, index: str, level: int = 0, parent_index: str | None = None) -> list[dict]:
    component_type = map_shape_type(shape)
    text = extract_text(shape)
    width_cm = emu_to_cm(shape.width)
    height_cm = emu_to_cm(shape.height)
    component = {
        "index": index,
        "parent_index": parent_index,
        "level": level,
        "name": shape.name,
        "type": component_type,
        "role": infer_role(component_type, text),
        "text": text,
        "is_text_editable": bool(getattr(shape, "has_text_frame", False)),
        "max_chars": estimate_max_chars(width_cm, height_cm, component_type),
        "position_cm": {
            "left": emu_to_cm(shape.left),
            "top": emu_to_cm(shape.top),
            "width": width_cm,
            "height": height_cm,
        },
    }
    rh = refine_heading_cap_category(shape, component_type)
    if rh:
        component["heading_cap_type"] = rh

    if getattr(shape, "has_text_frame", False) and component_type != "table":
        try:
            tr = extract_text_frame_style_runs(shape.text_frame)
            if tr:
                component["text_runs"] = tr
        except (AttributeError, TypeError, ValueError):
            pass

    if component_type == "group":
        child_shapes = list(shape.shapes)
        component["child_count"] = len(child_shapes)
        parsed_components = [component]
        for child_idx, child_shape in enumerate(child_shapes, start=1):
            child_index = f"{index}.{child_idx}"
            parsed_components.extend(
                parse_shape_component(
                    child_shape,
                    index=child_index,
                    level=level + 1,
                    parent_index=index,
                )
            )
        return parsed_components

    if component_type == "table":
        flat, nrows, ncols = flatten_table_text(shape)
        component["text"] = flat
        component["is_text_editable"] = False
        component["role"] = infer_role("table", flat)
        if nrows > 0 and ncols > 0:
            component["table_rows"] = nrows
            component["table_cols"] = ncols
        out: list[dict] = [component]
        if not getattr(shape, "has_table", False):
            return out
        try:
            tbl = shape.table
            nrows = len(tbl.rows)
            ncols = len(tbl.columns)
            cw = width_cm / max(ncols, 1)
            ch = height_cm / max(nrows, 1)
            for r in range(nrows):
                for c in range(ncols):
                    cell = tbl.cell(r, c)
                    ct = extract_text_from_cell(cell)
                    cell_index = f"{index}.t.{r + 1}.{c + 1}"
                    mxc = estimate_max_chars(cw, ch, "text")
                    cell_dict: dict[str, Any] = {
                        "index": cell_index,
                        "parent_index": index,
                        "level": level + 1,
                        "name": f"{shape.name or '表格'}-R{r + 1}C{c + 1}",
                        "type": "table_cell",
                        "role": infer_role("table_cell", ct),
                        "text": ct,
                        "is_text_editable": True,
                        "max_chars": mxc,
                        "table_row": r + 1,
                        "table_col": c + 1,
                        "position_cm": {
                            "left": emu_to_cm(shape.left),
                            "top": emu_to_cm(shape.top),
                            "width": cw,
                            "height": ch,
                        },
                    }
                    try:
                        tf = getattr(cell, "text_frame", None)
                        ctr = extract_text_frame_style_runs(tf)
                        if ctr:
                            cell_dict["text_runs"] = ctr
                    except (AttributeError, TypeError, ValueError):
                        pass
                    out.append(cell_dict)
        except (AttributeError, TypeError, ValueError):
            pass
        return out

    return [component]


def parse_pptx(file_path: Path) -> dict:
    prs = Presentation(str(file_path))
    result = {
        "file_name": file_path.name,
        "slide_count": len(prs.slides),
        "slides": [],
    }

    for idx, slide in enumerate(prs.slides, start=1):
        slide_info = {
            "slide_index": idx,
            "top_level_component_count": len(slide.shapes),
            "component_count": 0,
            "page_type": "unknown",
            "page_type_label": PAGE_TYPE_LABELS["unknown"],
            "page_type_reason": "",
            "components": [],
        }

        for s_idx, shape in enumerate(slide.shapes, start=1):
            slide_info["components"].extend(parse_shape_component(shape, index=str(s_idx)))

        slide_info["component_count"] = len(slide_info["components"])

        result["slides"].append(slide_info)

    return result


def iter_slide_shapes_with_index(slide):
    acc: list[tuple[str, object]] = []

    def visit(shape, idx: str) -> None:
        acc.append((idx, shape))
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for c_i, child in enumerate(shape.shapes, start=1):
                visit(child, f"{idx}.{c_i}")

    for i, sh in enumerate(slide.shapes, start=1):
        visit(sh, str(i))
    return acc


def get_shape_at_index(slide, index_path: str):
    """按与 iter_slide_shapes_with_index 相同的 1-based 路径解析形状（如 '2'、'2.3.1'）。"""
    parts = [p for p in (index_path or "").split(".") if p]
    if not parts:
        return None
    try:
        cur = slide.shapes[int(parts[0]) - 1]
        for p in parts[1:]:
            if not hasattr(cur, "shapes"):
                return None
            cur = cur.shapes[int(p) - 1]
        return cur
    except (IndexError, ValueError, TypeError, AttributeError):
        return None


def parse_table_cell_component_index(index_key: str) -> tuple[str, int, int] | None:
    """解析组件 index，如 '3.t.1.2' -> ('3', 1, 2)；'2.4.t.1.1' -> ('2.4', 1, 1)。"""
    key = str(index_key or "")
    if ".t." not in key:
        return None
    pos = key.index(".t.")
    prefix = key[:pos]
    tail = key[pos + 3 :]
    if not prefix:
        return None
    seg = tail.split(".")
    if len(seg) != 2:
        return None
    try:
        return prefix, int(seg[0]), int(seg[1])
    except ValueError:
        return None


def _primary_run_index(paragraph) -> int | None:
    runs = list(paragraph.runs)
    if not runs:
        return None
    best_i = 0
    best_len = -1
    for i, r in enumerate(runs):
        L = len((r.text or "").strip())
        if L > best_len:
            best_len = L
            best_i = i
    return best_i


def set_text_frame_plain(tf, text: str) -> None:
    if tf is None:
        return
    new_text = " ".join((text or "").replace("\r\n", "\n").split())
    try:
        if not tf.paragraphs:
            tf.text = new_text
            return
        p0 = tf.paragraphs[0]
        runs = list(p0.runs)
        if not runs:
            p0.text = new_text
        else:
            pi = _primary_run_index(p0)
            if pi is None:
                p0.text = new_text
            else:
                runs[pi].text = new_text
                for i, r in enumerate(runs):
                    if i != pi:
                        r.text = ""
        for p in tf.paragraphs[1:]:
            for r in p.runs:
                r.text = ""
    except Exception:  # noqa: BLE001
        tf.text = new_text


def set_shape_plain_text(shape, text: str) -> None:
    if not getattr(shape, "has_text_frame", False):
        return
    set_text_frame_plain(shape.text_frame, text)


def set_table_cell_plain_text(shape, row_1based: int, col_1based: int, text: str) -> None:
    if not getattr(shape, "has_table", False):
        return
    try:
        tbl = shape.table
        r, c = row_1based - 1, col_1based - 1
        n_r, n_c = len(tbl.rows), len(tbl.columns)
        if r < 0 or c < 0 or r >= n_r or c >= n_c:
            return
        cell = tbl.cell(r, c)
        tf = getattr(cell, "text_frame", None)
        if tf is not None:
            set_text_frame_plain(tf, text)
    except (AttributeError, TypeError, ValueError):
        pass


def _norm_slide_indices_list(slides: object) -> list[int]:
    seen: set[int] = set()
    ordered: list[int] = []
    if not isinstance(slides, list):
        return ordered
    for x in slides:
        try:
            n = int(x)
        except (TypeError, ValueError):
            continue
        if n not in seen:
            seen.add(n)
            ordered.append(n)
    return ordered


def _slides_by_index_map(parsed: dict) -> dict[int, dict[str, Any]]:
    out: dict[int, dict[str, Any]] = {}
    for s in parsed.get("slides") or []:
        if not isinstance(s, dict):
            continue
        try:
            idx = int(s.get("slide_index"))
        except (TypeError, ValueError):
            continue
        out[idx] = s
    return out


def _content_slide_indices_in_chapter(parsed: dict, chapter_slide_indices: list[int]) -> list[int]:
    """章节块内、且解析为正文页 (content) 的幻灯片页码，顺序与 chapter_slide_indices 一致。"""
    by_idx = _slides_by_index_map(parsed)
    out: list[int] = []
    for si in chapter_slide_indices:
        slide = by_idx.get(si)
        if not slide:
            continue
        if str(slide.get("page_type") or "").strip() == "content":
            out.append(si)
    return out


def _find_largest_picture_on_slide(slide):
    """当前页（含组合内）面积最大的图片形状。"""
    best = None
    best_area = -1
    for _idx, sh in iter_slide_shapes_with_index(slide):
        try:
            if sh.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            area = int(sh.width) * int(sh.height)
            if area > best_area:
                best_area = area
                best = sh
        except (AttributeError, TypeError, ValueError):
            continue
    return best


def _replace_picture_shape_preserving_z_order(shape, slide, image_path: Path) -> bool:
    """用磁盘图片替换已有图片占位，尽量保持位置、大小与在 spTree 中的顺序。"""
    if shape is None or shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
        return False
    path = Path(image_path)
    if not path.is_file():
        return False
    try:
        left, top, width, height = shape.left, shape.top, shape.width, shape.height
        sp = shape._element
        parent = sp.getparent()
        idx = list(parent).index(sp)
        parent.remove(sp)
        pic_shape = slide.shapes.add_picture(str(path), left, top, width=width, height=height)
        new_el = pic_shape._element
        new_par = new_el.getparent()
        new_par.remove(new_el)
        parent.insert(idx, new_el)
        return True
    except (AttributeError, TypeError, ValueError, OSError):
        return False


def apply_chapter_reference_screenshots(
    prs: Presentation,
    parsed: dict | None,
    chapter_ref: dict | None,
    task_id: str | None,
) -> None:
    """
    将各「章」Tab 上传的截图按顺序替换该章内正文页 (page_type=content) 上面积最大的图片。
    截图少于正文页时只替换前几页；不调用大模型。
    """
    if not isinstance(parsed, dict) or not isinstance(chapter_ref, dict):
        return
    tid = (task_id or "").strip()
    if not tid:
        return
    slots = chapter_ref.get("slots")
    if not isinstance(slots, list):
        return
    slot_rows = ppt_reference_slot_rows(parsed)
    if not slot_rows:
        return
    n = min(len(slot_rows), len(slots))
    by_parsed = _slides_by_index_map(parsed)

    for i in range(n):
        row = slot_rows[i]
        if str(row.get("kind") or "") != "chapter":
            continue
        slot = slots[i]
        if not isinstance(slot, dict):
            continue
        shots = slot.get("screenshots")
        if not isinstance(shots, list) or not shots:
            continue
        paths: list[Path] = []
        for sh in shots:
            if not isinstance(sh, dict):
                continue
            fn = str(sh.get("storedFilename") or "").strip()
            p = image_file_path(tid, fn) if fn else None
            if p is not None and p.is_file():
                paths.append(p)
        if not paths:
            continue
        chapter_pages = _norm_slide_indices_list(row.get("slides"))
        content_pages = _content_slide_indices_in_chapter(parsed, chapter_pages)
        path_i = 0
        for si in content_pages:
            if path_i >= len(paths):
                break
            if si < 1 or si > len(prs.slides):
                continue
            if not by_parsed.get(si):
                continue
            slide = prs.slides[si - 1]
            pic_shape = _find_largest_picture_on_slide(slide)
            if pic_shape is None:
                continue
            if _replace_picture_shape_preserving_z_order(pic_shape, slide, paths[path_i]):
                path_i += 1


def apply_generation_to_presentation(
    prs: Presentation,
    generated: dict,
    *,
    parsed: dict | None = None,
    chapter_ref: dict | None = None,
    task_id: str | None = None,
) -> None:
    for slide_payload in generated.get("slides", []):
        si = slide_payload.get("slide_index")
        if isinstance(si, str) and si.strip().isdigit():
            si = int(si)
        if not isinstance(si, int) or si < 1 or si > len(prs.slides):
            continue
        slide = prs.slides[si - 1]
        by_index = {
            str(c["index"]): str(c.get("generated_text") or "").strip()
            for c in slide_payload.get("components", [])
        }
        shape_by_idx = {idx: sh for idx, sh in iter_slide_shapes_with_index(slide)}
        for comp_index, text in by_index.items():
            if not text:
                continue
            parsed_cell = parse_table_cell_component_index(comp_index)
            if parsed_cell:
                prefix, row_i, col_i = parsed_cell
                tshape = get_shape_at_index(slide, prefix)
                if tshape is not None:
                    set_table_cell_plain_text(tshape, row_i, col_i, text)
                continue
            shape = shape_by_idx.get(comp_index)
            if shape is not None:
                set_shape_plain_text(shape, text)

    apply_chapter_reference_screenshots(prs, parsed, chapter_ref, task_id)

"""REST 风格 API：解析、生成、导出。"""
from __future__ import annotations

import io
import json
import mimetypes
import threading
import uuid
from typing import Any
from pathlib import Path

from flask import Blueprint, jsonify, request, send_file
from pptx import Presentation
from werkzeug.utils import secure_filename

from ppt_report import config
from ppt_report import state
from ppt_report.models import db as db_mod
from ppt_report.services.async_jobs import (
    _create_generate_job_record,
    _create_parse_job_record,
    _generate_job_worker,
    _generate_word_job_worker,
    _parse_job_worker,
    _word_parse_worker,
    snapshot_generate_job,
    snapshot_parse_job,
)
from ppt_report.services.generate_pipeline import merge_extra_from_upload, run_generate, validate_generate_prerequisites
from ppt_report.services.chapter_ref_images import (
    delete_chapter_ref_image,
    image_file_path,
    is_safe_task_id,
    save_chapter_ref_image,
)
from ppt_report.services.chapter_reference_resolve import resolve_chapter_reference
from ppt_report.services.page_types import compute_chapter_selection_groups
from ppt_report.services.presentation_cache import (
    bump_parsed_file_name_in_cache,
    get_parsed_from_cache,
    is_word_stored_payload,
    purge_presentation,
    resolve_template_path,
)
from ppt_report.services.pptx_document import apply_generation_to_presentation
from ppt_report.services.student_guidance_ai import generate_student_guidance
from ppt_report.services.student_import_ai import run_smart_import
from ppt_report.services.filled_export_cache import resolve_filled_export_path
from ppt_report.services.word_generation import fill_word_table_for_student
from ppt_report.utils.files import allowed_file

api_bp = Blueprint("api", __name__, url_prefix="/api")
WORD_REPORT_TEMPLATE_CODE = "word_table_fill"


def _is_word_report_type(template_id: str) -> bool:
    tid = (template_id or "").strip()
    if not tid:
        return False
    item = db_mod.get_chapter_template(tid)
    if not isinstance(item, dict):
        return False
    return str(item.get("templateCode") or "").strip() == WORD_REPORT_TEMPLATE_CODE


@api_bp.get("/presentations")
def api_presentations_list():
    return jsonify(
        {
            "ok": True,
            "items": db_mod.list_presentation_summaries(),
            "db_enabled": db_mod.db_enabled(),
        },
    )


@api_bp.get("/student-data")
def api_student_data_list():
    q = (request.args.get("q") or "").strip()
    items = db_mod.list_student_records(query=q)
    return jsonify({"ok": True, "items": items, "count": len(items)})


@api_bp.post("/student-data/ai-guidance")
def api_student_data_ai_guidance():
    """根据基础信息、学习画像、课时数据等生成成长指导四维文案。"""
    body: dict[str, Any] = request.get_json(silent=True) or {}
    profile = body.get("profile")
    if not isinstance(profile, dict):
        profile = {}
    content = body.get("content")
    content_s = (content if isinstance(content, str) else "").strip()
    try:
        guidance = generate_student_guidance(profile, content_s)
    except RuntimeError as exc:
        return jsonify({"ok": False, "error": str(exc)}), 400
    except Exception:
        return jsonify({"ok": False, "error": "生成失败，请稍后重试。"}), 500
    return jsonify({"ok": True, "guidance": guidance})


@api_bp.get("/student-data/<record_id>")
def api_student_data_get(record_id: str):
    item = db_mod.get_student_record(record_id)
    if not item:
        return jsonify({"ok": False, "error": "未找到该条学生数据。"}), 404
    return jsonify({"ok": True, "item": item})


@api_bp.post("/student-data")
def api_student_data_create():
    body: dict[str, Any] = request.get_json(silent=True) or {}
    try:
        record_id, item = db_mod.save_student_record(body)
    except ValueError as exc:
        return jsonify({"ok": False, "error": str(exc)}), 400
    except Exception:
        return jsonify({"ok": False, "error": "保存失败，请稍后重试。"}), 500
    return jsonify({"ok": True, "id": record_id, "item": item})


@api_bp.put("/student-data/<record_id>")
def api_student_data_update(record_id: str):
    body: dict[str, Any] = request.get_json(silent=True) or {}
    body["id"] = (record_id or "").strip()
    try:
        saved_id, item = db_mod.save_student_record(body)
    except ValueError as exc:
        return jsonify({"ok": False, "error": str(exc)}), 400
    except Exception:
        return jsonify({"ok": False, "error": "保存失败，请稍后重试。"}), 500
    return jsonify({"ok": True, "id": saved_id, "item": item})


@api_bp.delete("/student-data/<record_id>")
def api_student_data_delete(record_id: str):
    ok = db_mod.delete_student_record(record_id)
    if not ok:
        return jsonify({"ok": False, "error": "未找到该条学生数据。"}), 404
    return jsonify({"ok": True})


@api_bp.post("/student-data/import-ai")
def api_student_data_import_ai():
    """上传 CSV / xlsx，由大模型将行映射为系统档案字段后逐条入库（服务端按批调用模型）。"""
    if not db_mod.db_enabled():
        return jsonify({"ok": False, "error": "数据库未启用。"}), 400
    up = request.files.get("file")
    if not up or not (up.filename or "").strip():
        return jsonify({"ok": False, "error": "请选择要上传的 CSV 或 Excel 文件。"}), 400
    raw = up.read()
    if not raw:
        return jsonify({"ok": False, "error": "文件为空。"}), 400
    try:
        result = run_smart_import(raw, up.filename)
    except ValueError as exc:
        return jsonify({"ok": False, "error": str(exc)}), 400
    except RuntimeError as exc:
        return jsonify({"ok": False, "error": str(exc)}), 400
    except Exception:
        return jsonify({"ok": False, "error": "智能导入处理失败，请稍后重试。"}), 500
    return jsonify({"ok": True, **result})


@api_bp.get("/chapter-templates")
def api_chapter_templates_list():
    q = (request.args.get("q") or "").strip()
    items = db_mod.list_chapter_templates(query=q)
    return jsonify({"ok": True, "items": items, "count": len(items)})


@api_bp.get("/chapter-templates/<template_id>")
def api_chapter_templates_get(template_id: str):
    item = db_mod.get_chapter_template(template_id)
    if not item:
        return jsonify({"ok": False, "error": "未找到该模板。"}), 404
    return jsonify({"ok": True, "item": item})


@api_bp.post("/chapter-templates")
def api_chapter_templates_create():
    body: dict[str, Any] = request.get_json(silent=True) or {}
    try:
        template_id, item = db_mod.save_chapter_template(body)
    except ValueError as exc:
        return jsonify({"ok": False, "error": str(exc)}), 400
    except Exception:
        return jsonify({"ok": False, "error": "保存失败，请稍后重试。"}), 500
    return jsonify({"ok": True, "id": template_id, "item": item})


@api_bp.put("/chapter-templates/<template_id>")
def api_chapter_templates_update(template_id: str):
    body: dict[str, Any] = request.get_json(silent=True) or {}
    body["id"] = (template_id or "").strip()
    try:
        saved_id, item = db_mod.save_chapter_template(body)
    except KeyError:
        return jsonify({"ok": False, "error": "未找到该模板。"}), 404
    except ValueError as exc:
        return jsonify({"ok": False, "error": str(exc)}), 400
    except Exception:
        return jsonify({"ok": False, "error": "保存失败，请稍后重试。"}), 500
    return jsonify({"ok": True, "id": saved_id, "item": item})


@api_bp.delete("/chapter-templates/<template_id>")
def api_chapter_templates_delete(template_id: str):
    ok = db_mod.delete_chapter_template(template_id)
    if not ok:
        return jsonify({"ok": False, "error": "未找到该模板。"}), 404
    return jsonify({"ok": True})


@api_bp.post("/resolve-chapter-reference")
def api_resolve_chapter_reference():
    """
    按章节模板顺序写入各「章」块的标题，并由大模型将学生字段分配到各章。
    请求 JSON：task_id, chapter_template_id, student_data_id；可选 use_llm（默认 true）。
    """
    body = request.get_json(silent=True) or {}
    task_id = (body.get("task_id") or "").strip()
    chapter_template_id = (body.get("chapter_template_id") or "").strip()
    student_data_id = (body.get("student_data_id") or "").strip()
    parsed_chk = get_parsed_from_cache(task_id)
    if not parsed_chk:
        return jsonify({"ok": False, "error": "未找到该解析记录。"}), 404
    if is_word_stored_payload(parsed_chk if isinstance(parsed_chk, dict) else None):
        return jsonify({"ok": False, "error": "Word 文档不能用于章节解析。"}), 400
    use_llm = body.get("use_llm")
    if use_llm is None:
        use_llm = True
    elif isinstance(use_llm, str):
        use_llm = use_llm.lower() not in ("0", "false", "no")
    else:
        use_llm = bool(use_llm)
    try:
        data = resolve_chapter_reference(
            task_id,
            chapter_template_id,
            student_data_id,
            use_llm=use_llm,
        )
    except ValueError as e:
        return jsonify({"ok": False, "error": str(e)}), 400
    except RuntimeError as e:
        return jsonify({"ok": False, "error": str(e)}), 503
    except Exception:
        return jsonify({"ok": False, "error": "解析失败，请稍后重试。"}), 500
    return jsonify({"ok": True, "data": data})


@api_bp.post("/chapter-ref-screenshot")
def api_chapter_ref_screenshot_upload():
    """上传某一章的参考截图（按 task_id 分目录存储）。"""
    task_id = (request.form.get("task_id") or "").strip()
    if not get_parsed_from_cache(task_id):
        return jsonify({"ok": False, "error": "未找到该 PPT 解析记录，无法上传。"}), 404
    file = request.files.get("file")
    item, err = save_chapter_ref_image(task_id, file)
    if err:
        code = 413 if "不能超过" in err else 400
        return jsonify({"ok": False, "error": err}), code
    return jsonify({"ok": True, "item": item})


@api_bp.delete("/chapter-ref-screenshot/<task_id>/<filename>")
def api_chapter_ref_screenshot_delete(task_id: str, filename: str):
    """删除已上传的参考截图。"""
    ok = delete_chapter_ref_image(task_id.strip(), filename.strip())
    if not ok:
        return jsonify({"ok": False, "error": "未找到该文件。"}), 404
    return jsonify({"ok": True})


@api_bp.get("/chapter-ref-images/<task_id>/<filename>")
def api_chapter_ref_image_get(task_id: str, filename: str):
    """读取章节参考截图（供预览与后续填入 PPT）。"""
    path = image_file_path(task_id.strip(), filename.strip())
    if not path:
        return jsonify({"ok": False, "error": "未找到文件。"}), 404
    mt = mimetypes.guess_type(path.name)[0] or "application/octet-stream"
    return send_file(path, mimetype=mt)


@api_bp.delete("/presentations/<task_id>")
def api_presentation_delete(task_id: str):
    """删除解析记录、模板文件、内存缓存及该任务的章节参考截图。"""
    ok, err = purge_presentation(task_id)
    if not ok:
        code = 404 if err and "未找到" in (err or "") else 400
        return jsonify({"ok": False, "error": err or "删除失败"}), code
    return jsonify({"ok": True})


@api_bp.patch("/presentations/<task_id>")
def api_presentation_rename(task_id: str):
    """修改已解析模板在数据库中的显示文件名（不改动磁盘 .pptx 文件名）。"""
    if not db_mod.db_enabled():
        return jsonify({"ok": False, "error": "数据库未启用，无法保存。"}), 503
    body = request.get_json(silent=True) or {}
    file_name = body.get("file_name")
    if file_name is None:
        file_name = body.get("name")
    if file_name is None:
        return jsonify({"ok": False, "error": "请提供 file_name。"}), 400
    fn = str(file_name).strip()
    if not fn:
        return jsonify({"ok": False, "error": "file_name 不能为空。"}), 400
    ok, err = db_mod.update_presentation_file_name(task_id, fn)
    if not ok:
        code = 404 if err and "未找到" in err else 400
        return jsonify({"ok": False, "error": err or "更新失败"}), code
    tid = (task_id or "").strip()
    row_name = fn[:512]
    bump_parsed_file_name_in_cache(tid, row_name)
    return jsonify({"ok": True, "task_id": tid, "file_name": row_name})


@api_bp.get("/presentations/<task_id>")
def api_presentation_for_generate(task_id: str):
    tid = (task_id or "").strip()
    parsed = get_parsed_from_cache(tid)
    if not parsed:
        return jsonify({"ok": False, "error": "未找到该解析记录。"}), 404
    tpl = resolve_template_path(tid)
    return jsonify(
        {
            "ok": True,
            "task_id": tid,
            "file_name": str(parsed.get("file_name") or ""),
            "slide_count": int(parsed.get("slide_count") or 0),
            "chapter_groups": compute_chapter_selection_groups(parsed),
            "has_template": bool(tpl and tpl.is_file()),
        },
    )


@api_bp.get("/generation_history")
def api_generation_history_list():
    query = (request.args.get("q") or "").strip()
    return jsonify(
        {
            "ok": True,
            "db_enabled": db_mod.db_enabled(),
            "items": db_mod.list_generation_history_summaries(query=query),
        },
    )


@api_bp.get("/generation_history/<record_id>")
def api_generation_history_get(record_id: str):
    rec = db_mod.get_generation_history(record_id)
    if not rec:
        return jsonify({"ok": False, "error": "记录不存在。"}), 404
    return jsonify({"ok": True, "record": rec})


@api_bp.delete("/generation_history/<record_id>")
def api_generation_history_delete(record_id: str):
    if not db_mod.db_enabled():
        return jsonify({"ok": False, "error": "数据库未启用。"}), 503
    if not db_mod.delete_generation_history(record_id):
        return jsonify({"ok": False, "error": "记录不存在。"}), 404
    return jsonify({"ok": True})


@api_bp.post("/generation_history/cleanup")
def api_generation_history_cleanup():
    if not db_mod.db_enabled():
        return jsonify({"ok": False, "error": "数据库未启用。"}), 503
    data = request.get_json(silent=True) or {}
    try:
        days = int(data.get("days") or 0)
    except (TypeError, ValueError):
        days = 0
    if days not in (3, 7, 15):
        return jsonify({"ok": False, "error": "仅支持按 3/7/15 天清理。"}), 400
    removed = db_mod.cleanup_expired_generation_history(retention_days=days)
    return jsonify({"ok": True, "removed": int(removed), "days": days})


@api_bp.get("/generation_history/<record_id>/download")
def api_generation_history_download(record_id: str):
    """
    历史下载：
    1) 优先返回 uploads/filled_exports/{record_id}.pptx
    2) 若缓存不存在，回退到「按 task_id + result 即时导出」
    3) 若两者均不可用，给出明确提示
    """
    rec = db_mod.get_generation_history(record_id)
    if not rec:
        return jsonify({"ok": False, "error": "记录不存在，可能已被删除。"}), 404

    result_obj = rec.get("result") if isinstance(rec.get("result"), dict) else {}
    output_kind = "docx" if str(result_obj.get("output_kind") or "").strip() == "docx" else "pptx"

    cached = resolve_filled_export_path(record_id, output_kind)
    if cached:
        is_docx = cached.suffix.lower() == ".docx"
        return send_file(
            cached,
            as_attachment=True,
            download_name=f"history_{record_id}_filled{'.docx' if is_docx else '.pptx'}",
            mimetype=(
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                if is_docx
                else "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            ),
        )

    if output_kind == "docx":
        fill = result_obj.get("word_fill_summary") if isinstance(result_obj.get("word_fill_summary"), dict) else {}
        task_id = str(fill.get("task_id") or rec.get("taskId") or "").strip()
        student_data_id = str(fill.get("student_data_id") or "").strip()
        if not task_id or not student_data_id:
            return (
                jsonify(
                    {
                        "ok": False,
                        "error": "历史缓存文件不存在，且记录缺少 Word 回填所需数据，无法下载。请重新生成。",
                    },
                ),
                410,
            )
        out = config.FILLED_EXPORT_DIR / f"{record_id}.docx"
        try:
            fill_word_table_for_student(task_id, student_data_id, output_path=out)
        except Exception as exc:  # noqa: BLE001
            return jsonify({"ok": False, "error": str(exc)}), 410
        return send_file(
            out,
            as_attachment=True,
            download_name=f"history_{record_id}_filled.docx",
            mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        )

    task_id = str(rec.get("taskId") or "").strip()
    generated = rec.get("result")
    if not task_id or not isinstance(generated, dict):
        return (
            jsonify(
                {
                    "ok": False,
                    "error": "历史缓存文件不存在，且记录缺少导出所需数据，无法下载。请重新生成。",
                },
            ),
            410,
        )

    template_path = resolve_template_path(task_id)
    if not template_path or not template_path.is_file():
        return (
            jsonify(
                {
                    "ok": False,
                    "error": "历史缓存文件不存在，且原始 PPT 模板已缺失（可能被清理或未持久化），无法下载。请重新上传并生成。",
                },
            ),
            410,
        )

    prs = Presentation(str(template_path))
    parsed_export = get_parsed_from_cache(task_id)
    cr = state.LAST_CHAPTER_REF.get(task_id)
    if not isinstance(cr, dict):
        cr = None
    apply_generation_to_presentation(
        prs,
        generated,
        parsed=parsed_export,
        chapter_ref=cr,
        task_id=task_id,
    )
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    meta = parsed_export or {}
    raw_name = meta.get("file_name") or "presentation.pptx"
    stem = Path(raw_name).stem
    return send_file(
        buf,
        as_attachment=True,
        download_name=f"{stem}_filled.pptx",
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )


@api_bp.post("/parse_start")
def api_parse_start():
    file = request.files.get("ppt_file")
    if not file or not file.filename:
        return jsonify({"ok": False, "error": "请先选择一个文件。"}), 400
    if not allowed_file(file.filename):
        return jsonify({"ok": False, "error": "只支持 .pptx、.ppt 或 .docx 文件。"}), 400
    lower = file.filename.lower()
    is_docx = lower.endswith(".docx")
    if lower.endswith(".ppt") and not lower.endswith(".pptx"):
        return jsonify({"ok": False, "error": "当前仅支持 .pptx 解析，请先将 .ppt 转为 .pptx。"}), 400

    replace_raw = (request.form.get("replace_task_id") or "").strip()
    replace_tid = replace_raw if replace_raw and is_safe_task_id(replace_raw) else None
    if replace_tid:
        prev = get_parsed_from_cache(replace_tid)
        prev_word = is_word_stored_payload(prev) if isinstance(prev, dict) else False
        if prev is None:
            return jsonify({"ok": False, "error": "未找到要替换的模板记录，可能已删除。"}), 400
        if prev_word and not is_docx:
            return jsonify({"ok": False, "error": "该条目为 Word 模板，请仅上传 .docx 替换。"}), 400
        if not prev_word and is_docx:
            return jsonify({"ok": False, "error": "该条目为已解析的 PPT，请上传 .pptx 替换。"}), 400

    orig_name = secure_filename(file.filename) or ("document.docx" if is_docx else "presentation.pptx")
    suffix = ".docx" if is_docx else ".pptx"
    tmp_path = config.UPLOAD_DIR / f"_parse_tmp_{uuid.uuid4().hex}{suffix}"
    file.save(str(tmp_path))

    job_poll_id = _create_parse_job_record()
    if is_docx:
        threading.Thread(
            target=_word_parse_worker,
            args=(job_poll_id, tmp_path, orig_name, replace_tid),
            daemon=True,
        ).start()
    else:
        threading.Thread(
            target=_parse_job_worker,
            args=(job_poll_id, tmp_path, orig_name, replace_tid),
            daemon=True,
        ).start()
    return jsonify({"ok": True, "job_id": job_poll_id})


@api_bp.get("/parse_status/<job_id>")
def api_parse_status(job_id: str):
    snap = snapshot_parse_job(job_id)
    if not snap:
        return jsonify({"ok": False, "error": "任务不存在或已过期。"}), 404
    return jsonify(snap)


@api_bp.post("/generate_start")
def api_generate_start():
    task_id = request.form.get("task_id")
    topic = (request.form.get("topic") or "").strip()
    chapter_template_id = (request.form.get("chapter_template_id") or "").strip()
    student_data_id = (request.form.get("student_data_id") or "").strip()
    selected_slides = sorted(
        {
            int(x)
            for x in request.form.getlist("selected_slides")
            if str(x).isdigit()
        }
    )
    chapter_ref: dict | None = None
    raw_ref = (request.form.get("chapter_ref_json") or "").strip()
    if raw_ref:
        try:
            loaded = json.loads(raw_ref)
            if isinstance(loaded, dict):
                chapter_ref = loaded
        except (json.JSONDecodeError, TypeError, ValueError):
            chapter_ref = None
    err, merged_extra = merge_extra_from_upload("", None)
    if err:
        return jsonify({"ok": False, "error": err}), 400
    if not chapter_template_id:
        return jsonify({"ok": False, "error": "请选择报告类型。"}), 400
    if _is_word_report_type(chapter_template_id):
        parsed = get_parsed_from_cache(task_id)
        if not parsed:
            return jsonify({"ok": False, "error": "未找到模板记录，请重新选择。"}), 400
        if not is_word_stored_payload(parsed if isinstance(parsed, dict) else None):
            return jsonify({"ok": False, "error": "Word 报告类型仅支持 Word 模板。"}), 400
        if not student_data_id:
            return jsonify({"ok": False, "error": "请选择学生数据。"}), 400
        job_id = _create_generate_job_record(task_id)
        threading.Thread(
            target=_generate_word_job_worker,
            args=(job_id, task_id, topic, merged_extra, selected_slides),
            kwargs={"student_data_id": student_data_id},
            daemon=True,
        ).start()
        return jsonify({"ok": True, "job_id": job_id})
    err = validate_generate_prerequisites(
        task_id, topic, merged_extra, selected_slides, chapter_ref
    )
    if err:
        return jsonify({"ok": False, "error": err}), 400
    job_id = _create_generate_job_record(task_id)
    threading.Thread(
        target=_generate_job_worker,
        args=(job_id, task_id, topic, merged_extra, selected_slides, chapter_ref),
        daemon=True,
    ).start()
    return jsonify({"ok": True, "job_id": job_id})


@api_bp.get("/generate_status/<job_id>")
def api_generate_status(job_id: str):
    snap = snapshot_generate_job(job_id)
    if not snap:
        return jsonify({"ok": False, "error": "任务不存在或已过期。"}), 404
    return jsonify(snap)


@api_bp.post("/generate")
def api_generate():
    data = request.get_json(silent=True) or {}
    task_id = data.get("task_id")
    topic = data.get("topic") or ""
    extra_content = data.get("extra_content") or ""
    chapter_template_id = str(data.get("chapter_template_id") or "").strip()
    student_data_id = str(data.get("student_data_id") or "").strip()
    slides = data.get("selected_slides") or data.get("slides") or []
    selected_slides = sorted(
        {
            int(x)
            for x in slides
            if isinstance(x, int) or (isinstance(x, str) and x.strip().isdigit())
        }
    )

    chapter_ref = data.get("chapter_ref")
    if not isinstance(chapter_ref, dict):
        chapter_ref = None
    if chapter_ref is None:
        crj = data.get("chapter_ref_json")
        if isinstance(crj, str) and crj.strip():
            try:
                loaded = json.loads(crj)
                if isinstance(loaded, dict):
                    chapter_ref = loaded
            except (json.JSONDecodeError, TypeError, ValueError):
                chapter_ref = None

    if chapter_template_id and _is_word_report_type(chapter_template_id):
        try:
            summary = fill_word_table_for_student(str(task_id or "").strip(), student_data_id)
        except Exception as exc:  # noqa: BLE001
            return jsonify({"ok": False, "error": str(exc)}), 400
        return jsonify({"ok": True, "data": {"output_kind": "docx", "word_fill_summary": summary}})

    err, result, _ = run_generate(
        task_id, topic, extra_content, selected_slides, None, chapter_ref=chapter_ref
    )
    if err:
        return jsonify({"ok": False, "error": err}), 400
    return jsonify({"ok": True, "data": result})


@api_bp.post("/export")
def api_export():
    body = request.get_json(silent=True) or {}
    task_id = body.get("task_id")
    generated = body.get("data") or (state.LAST_GENERATION.get(task_id) if task_id else None)
    template_path = resolve_template_path(task_id) if task_id else None
    if not task_id:
        return jsonify({"ok": False, "error": "缺少 task_id。"}), 400
    if not template_path or not template_path.is_file():
        return jsonify({"ok": False, "error": "未找到模板文件。"}), 400
    if not generated:
        return jsonify({"ok": False, "error": "请先生成文档内容，或在 body 中传入 data。"}), 400
    prs = Presentation(str(template_path))
    parsed_export = get_parsed_from_cache(task_id)
    cr = body.get("chapter_ref")
    if not isinstance(cr, dict):
        crj = body.get("chapter_ref_json")
        if isinstance(crj, str) and crj.strip():
            try:
                loaded = json.loads(crj)
                cr = loaded if isinstance(loaded, dict) else None
            except (json.JSONDecodeError, TypeError, ValueError):
                cr = None
        if not isinstance(cr, dict):
            cr = state.LAST_CHAPTER_REF.get(task_id)
    if not isinstance(cr, dict):
        cr = None
    apply_generation_to_presentation(
        prs,
        generated,
        parsed=parsed_export,
        chapter_ref=cr,
        task_id=task_id,
    )
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    meta = parsed_export or {}
    raw_name = meta.get("file_name") or "presentation.pptx"
    stem = Path(raw_name).stem
    return send_file(
        buf,
        as_attachment=True,
        download_name=f"{stem}_filled.pptx",
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )


# 与历史路径一致：/export/<task_id> 挂在应用根上，不单挂载在 /api
export_bp = Blueprint("export", __name__)


@export_bp.get("/export/<task_id>")
def export_filled_pptx(task_id: str):
    template_path = resolve_template_path(task_id)
    generated = state.LAST_GENERATION.get(task_id)
    if not template_path or not template_path.is_file():
        return "未找到该任务的模板，请重新上传并解析。", 404
    if not generated:
        return "请先生成文档内容，再下载填充后的 PPT。", 400
    prs = Presentation(str(template_path))
    parsed_export = get_parsed_from_cache(task_id)
    cr = state.LAST_CHAPTER_REF.get(task_id)
    if not isinstance(cr, dict):
        cr = None
    apply_generation_to_presentation(
        prs,
        generated,
        parsed=parsed_export,
        chapter_ref=cr,
        task_id=task_id,
    )
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    meta = parsed_export or {}
    raw_name = meta.get("file_name") or "presentation.pptx"
    stem = Path(raw_name).stem
    return send_file(
        buf,
        as_attachment=True,
        download_name=f"{stem}_filled.pptx",
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )
